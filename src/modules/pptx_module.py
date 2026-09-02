"""
PowerPoint (.pptx) Document Module for DocumentConvertTool.
Converts PPTX slides to Markdown and Markdown to PPTX slides.
"""
import os
import re
import math
import zipfile
from src.core.base_module import BaseDocumentModule
from src.core.registry import ModuleRegistry


class PPTXModule(BaseDocumentModule):
    @property
    def name(self) -> str:
        return "PowerPoint"

    @property
    def file_extensions(self) -> list[str]:
        return [".pptx"]

    @property
    def required_dependencies(self) -> list[str]:
        return ["python-pptx", "Pillow"]

    # -------------------------------------------------------------------------
    # Private Helper Methods (Un-nested from load_to_markdown & save_from_markdown)
    # -------------------------------------------------------------------------

    @staticmethod
    def _is_smartart_shape(shape, mso_shape_type) -> bool:
        """Checks if shape is SmartArt or GraphicFrame containing diagram XML."""
        try:
            diag_enum = getattr(mso_shape_type, "DIAGRAM", None)
            if diag_enum is not None and getattr(shape, "shape_type", None) == diag_enum:
                return True
        except Exception:
            pass

        try:
            xml_str = getattr(shape.element, "xml", "")
            if "drawingml/2006/diagram" in xml_str or "dgm:" in xml_str:
                return True
        except Exception:
            pass
        return False

    @staticmethod
    def _is_numbered_paragraph(paragraph) -> bool:
        """Checks if a python-pptx paragraph uses PowerPoint auto-numbering (buAutoNum)."""
        try:
            pPr = getattr(getattr(paragraph, "_p", None), "pPr", None)
            if pPr is not None:
                xml_str = getattr(pPr, "xml", "")
                if "buAutoNum" in xml_str:
                    return True
                if hasattr(paragraph._p, "xpath") and paragraph._p.xpath("./a:pPr/a:buAutoNum"):
                    return True
        except Exception:
            pass
        return False

    @classmethod
    def _flatten_shapes(cls, shapes_iterable, mso_shape_type):
        """Recursively flattens group shapes and sorts them by top position."""
        flat = []
        sorted_items = sorted(shapes_iterable, key=lambda s: getattr(s, "top", 0) or 0)
        for s in sorted_items:
            if hasattr(s, "shape_type") and s.shape_type == mso_shape_type.GROUP:
                if hasattr(s, "shapes"):
                    flat.extend(cls._flatten_shapes(s.shapes, mso_shape_type))
                else:
                    flat.append(s)
            else:
                flat.append(s)
        return flat

    @staticmethod
    def _split_markdown_into_slide_blocks(content: str) -> list[str]:
        """
        Splits markdown into slide blocks.
        If explicit '\n\n---\n\n' dividers exist, split strictly by them.
        Otherwise, auto-segment by major Markdown headings (#, ##, ###) into separate slides.
        """
        if "\n\n---\n\n" in content:
            blocks = [s.strip() for s in content.split("\n\n---\n\n") if s.strip()]
            if blocks:
                return blocks

        lines = content.split("\n")
        blocks = []
        current = []
        in_code = False

        for l in lines:
            stripped = l.strip()
            if stripped.startswith("```"):
                in_code = not in_code
                current.append(l)
                continue

            is_notes_heading = bool(re.match(r"^#{1,3}\s*Notes:", stripped, re.IGNORECASE))
            is_heading = not in_code and bool(re.match(r"^#{1,3}\s+\S+", stripped)) and not is_notes_heading
            if is_heading and current:
                non_empty = [x for x in current if x.strip()]
                if len(non_empty) >= 2:
                    blocks.append("\n".join(current).strip())
                    current = []
            current.append(l)

        if current:
            blocks.append("\n".join(current).strip())

        return [b for b in blocks if b.strip()]

    @staticmethod
    def _get_visual_line_count(text: str, level: int = 0) -> int:
        """
        Estimates visual wrapped line count for a paragraph line based on character count,
        taking into account text frame width (~11.7 in) and level indentations.
        """
        clean = text.strip()
        if not clean:
            return 1
        chars_per_line = max(35, 75 - level * 5)
        return max(1, math.ceil(len(clean) / chars_per_line))

    @staticmethod
    def _parse_formatted_runs(paragraph, text: str, pt_cls, default_size=16, force_bold=False):
        """Parses inline bold/italic/code/underline/strikethrough markdown text and appends runs to PPTX paragraph."""
        from src.core.converters import parse_inline

        segments = parse_inline(text, bold=force_bold)
        for seg in segments:
            if not seg.text:
                continue
            run = paragraph.add_run()
            run.text = seg.text
            if seg.code:
                run.font.name = "Consolas"
                run.font.size = pt_cls(max(15, default_size))
            else:
                run.font.size = pt_cls(default_size)

            if force_bold or seg.bold:
                run.font.bold = True
            if seg.italic:
                run.font.italic = True
            if seg.underline or (seg.url and not seg.is_image):
                run.font.underline = True
            if seg.strike:
                try:
                    run.font.strike = True
                except Exception:
                    pass

            if seg.url and not seg.is_image:
                try:
                    run.hyperlink.address = seg.url
                except Exception as ex:
                    print(f"[DEBUG] Failed to set PPTX run hyperlink: {ex}")

    @staticmethod
    def _calc_title_geometry(title: str, inches_cls) -> tuple:
        """
        Calculates (body_top_offset, actual_title_height, title_font_size) based on title string length.
        Snaps title box height tightly to actual title content height and leaves a clean gap for body text.
        """
        if not title:
            return inches_cls(0.8), inches_cls(0.0), 32
        length = len(title)
        if length > 75:
            title_lines = 3
            font_sz = 20
        elif length > 40:
            title_lines = 2
            font_sz = 22
        else:
            title_lines = 1
            font_sz = 26

        title_height = inches_cls(0.42 * title_lines + 0.1)
        body_top_offset = inches_cls(0.5) + title_height + inches_cls(0.4)
        return body_top_offset, title_height, font_sz

    @classmethod
    def _create_continuation_slide(cls, prs, title_text, notes_text, blank_layout, inches_cls, pt_cls, pp_align_cls) -> tuple:
        """Helper to create unified continuation slide with snapped title height and left alignment."""
        if title_text:
            title_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(title_layout)
            cont_title = f"{title_text} (Cont.)"
            c_offset, actual_title_h, c_font_sz = cls._calc_title_geometry(cont_title, inches_cls)
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = cont_title
                try:
                    title_shape.left = inches_cls(0.8)
                    title_shape.top = inches_cls(0.5)
                    title_shape.width = inches_cls(11.7)
                    title_shape.height = actual_title_h
                    tf_t = title_shape.text_frame
                    tf_t.word_wrap = True
                    tf_t.margin_top = inches_cls(0.0)
                    tf_t.margin_bottom = inches_cls(0.0)
                    tf_t.margin_left = inches_cls(0.0)
                    tf_t.margin_right = inches_cls(0.0)
                    for p in tf_t.paragraphs:
                        if pp_align_cls is not None:
                            try:
                                p.alignment = pp_align_cls.LEFT
                            except Exception:
                                pass
                        for r in p.runs:
                            r.font.size = pt_cls(c_font_sz)
                except Exception:
                    pass
            top_offset = c_offset
            start_top_offset = c_offset
        else:
            slide = prs.slides.add_slide(blank_layout)
            top_offset = inches_cls(0.8)
            start_top_offset = inches_cls(0.8)
        if notes_text:
            try:
                slide.notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass
        return slide, top_offset, start_top_offset

    @classmethod
    def _ensure_space(cls, slide, top_offset, start_top_offset, max_slide_top, prs, title_text, notes_text, blank_layout, inches_cls, pt_cls, pp_align_cls, required_height=None) -> tuple:
        if required_height is None:
            required_height = inches_cls(1.0)
        if top_offset + required_height > max_slide_top and top_offset > start_top_offset:
            print(f"[DEBUG] Slide content overflow warning: top_offset={top_offset} + {required_height} exceeds slide height, auto-creating continuation slide")
            return cls._create_continuation_slide(prs, title_text, notes_text, blank_layout, inches_cls, pt_cls, pp_align_cls)
        return slide, top_offset, start_top_offset

    @staticmethod
    def _calc_row_h(row: list[str], col_char_capacities: list[int], inches_cls) -> float:
        max_lines = 1
        for c_idx, cell in enumerate(row):
            cap = col_char_capacities[c_idx] if c_idx < len(col_char_capacities) else 25
            clean_cell = re.sub(r"\*\*|\*|`", "", cell.replace("<br>", "\n")).strip()
            cell_lines = sum(max(1, math.ceil(len(l_str) / cap)) for l_str in clean_cell.split("\n"))
            max_lines = max(max_lines, cell_lines)
        return inches_cls(0.26 * max_lines + 0.18)

    # -------------------------------------------------------------------------
    # BaseDocumentModule Abstract Methods Implementation
    # -------------------------------------------------------------------------

    def load_to_markdown(self, file_path: str) -> str:
        """
        Extracts PowerPoint (.pptx) slides to clean Markdown text.
        Preserves titles (as ## headings), bullet levels, tables, pictures via MediaAssetManager,
        charts (as MD tables), and slide notes.
        Slides are separated strictly by '\\n\\n---\\n\\n'.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        if not zipfile.is_zipfile(file_path):
            return "⚠️ File Validation Error: Invalid file. Please select a valid PPTX document."

        try:
            import pptx
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
            from src.core.converters import wrap_text_style
            from src.services.media_asset_manager import MediaAssetManager

            asset_mgr = MediaAssetManager()
            prs = pptx.Presentation(file_path)
            slide_markdowns = []

            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_blocks = []
                all_shapes = self._flatten_shapes(slide.shapes, MSO_SHAPE_TYPE)

                for shape in all_shapes:
                    # 1. Slide Title Detection
                    if getattr(shape, "is_placeholder", False):
                        try:
                            p_type = shape.placeholder_format.type
                            if p_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                                if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                                    title_text = shape.text_frame.text.strip().replace("\n", " ")
                                    slide_blocks.append(f"## {title_text}")
                                    continue
                        except Exception:
                            pass

                    # 2. Table extraction
                    if shape.has_table:
                        try:
                            table = shape.table
                            table_rows = []
                            for row in table.rows:
                                cell_texts = [cell.text.replace("\n", "<br>").strip() for cell in row.cells]
                                table_rows.append(f"| {' | '.join(cell_texts)} |")

                            if table_rows:
                                col_count = len(table.columns)
                                sep = f"| {' | '.join(['---'] * col_count)} |"
                                table_md = "\n".join([table_rows[0], sep] + table_rows[1:])
                                slide_blocks.append(table_md)
                                continue
                        except Exception as e:
                            print(f"[DEBUG] Table extraction error in PPTX: {e}")

                    # 3. Chart extraction -> Markdown Table
                    if shape.has_chart:
                        try:
                            chart = shape.chart
                            plots = chart.plots
                            if plots:
                                categories = [str(c) for c in plots[0].categories]
                                series_list = list(chart.series)

                                if categories and series_list:
                                    headers = ["Category"] + [s.name or f"Series {i+1}" for i, s in enumerate(series_list)]
                                    header_line = f"| {' | '.join(headers)} |"
                                    sep_line = f"| {' | '.join(['---'] * len(headers))} |"

                                    data_lines = []
                                    for c_idx, cat in enumerate(categories):
                                        row_vals = [cat]
                                        for s in series_list:
                                            vals = list(s.values)
                                            val_str = str(vals[c_idx]) if c_idx < len(vals) else ""
                                            row_vals.append(val_str)
                                        data_lines.append(f"| {' | '.join(row_vals)} |")

                                    chart_title = ""
                                    try:
                                        if getattr(chart, "has_title", False) and hasattr(chart, "chart_title"):
                                            tf = getattr(chart.chart_title, "text_frame", None)
                                            if tf and hasattr(tf, "paragraphs"):
                                                extracted_t = "\n".join(p.text for p in tf.paragraphs if p.text and p.text.strip()).strip()
                                                chart_title = extracted_t if extracted_t else "Tiêu đề Biểu đồ"
                                            else:
                                                chart_title = "Tiêu đề Biểu đồ"
                                    except Exception:
                                        pass

                                    chart_type_name = "column"
                                    try:
                                        c_type = getattr(chart, "chart_type", None)
                                        if c_type is not None:
                                            c_type_str = str(c_type).lower()
                                            if "bar" in c_type_str:
                                                chart_type_name = "bar"
                                            elif "pie" in c_type_str:
                                                chart_type_name = "pie"
                                            elif "line" in c_type_str:
                                                chart_type_name = "line"
                                    except Exception:
                                        pass

                                    legend_pos = ""
                                    try:
                                        if getattr(chart, "has_legend", False) and hasattr(chart, "legend"):
                                            pos_enum = getattr(chart.legend, "position", None)
                                            if pos_enum is not None:
                                                pos_str = str(pos_enum).lower()
                                                if "right" in pos_str:
                                                    legend_pos = "right"
                                                elif "top" in pos_str:
                                                    legend_pos = "top"
                                                elif "left" in pos_str:
                                                    legend_pos = "left"
                                                elif "bottom" in pos_str:
                                                    legend_pos = "bottom"
                                    except Exception:
                                        pass

                                    title_attr = f' title="{chart_title}"' if chart_title else ""
                                    legend_attr = f' legend="{legend_pos}"' if legend_pos else ""
                                    chart_md = f"<!-- chart: {chart_type_name}{title_attr}{legend_attr} -->\n" + "\n".join([header_line, sep_line] + data_lines)
                                    slide_blocks.append(chart_md)
                                    continue
                        except Exception as e:
                            print(f"[DEBUG] Chart extraction error in PPTX: {e}")
                            slide_blocks.append("*[Unsupported: Chart]*")
                            continue

                    # 4. Picture extraction
                    if (hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE) or hasattr(shape, "image"):
                        try:
                            image_blob = shape.image.blob
                            ext = getattr(shape.image, "ext", "png") or "png"
                            img_filename = f"pptx_slide{slide_idx}_img_{id(shape)}.{ext}"
                            virtual_uri = asset_mgr.register_image(image_blob, img_filename)
                            slide_blocks.append(f"![image]({virtual_uri})")
                            continue
                        except Exception as e:
                            print(f"[DEBUG] Picture extraction error in PPTX: {e}")

                    # 5. SmartArt detection & fallback
                    if self._is_smartart_shape(shape, MSO_SHAPE_TYPE):
                        slide_blocks.append("*[Unsupported: SmartArt/Diagram]*")
                        continue

                    # 6. Text Frame / Paragraphs (Body text)
                    if shape.has_text_frame:
                        tf = shape.text_frame
                        para_lines = []
                        num_counters = {}
                        for paragraph in tf.paragraphs:
                            if not paragraph.text or not paragraph.text.strip():
                                continue

                            run_parts = []
                            for run in paragraph.runs:
                                if not run.text:
                                    continue
                                font = getattr(run, "font", None)
                                bold = getattr(font, "bold", None) if font else None
                                italic = getattr(font, "italic", None) if font else None
                                strike = getattr(font, "strike", None) if font else None
                                underline = getattr(font, "underline", None) if font else None

                                formatted = wrap_text_style(
                                    run.text,
                                    bold=bold,
                                    italic=italic,
                                    strike=strike,
                                    underline=underline
                                )
                                run_parts.append(formatted)

                            p_text = "".join(run_parts).strip()
                            if p_text:
                                lvl = getattr(paragraph, "level", 0) or 0
                                indent = "  " * lvl
                                is_num = self._is_numbered_paragraph(paragraph) or bool(re.match(r"^\d+\.\s+", p_text))
                                if is_num:
                                    num_counters[lvl] = num_counters.get(lvl, 0) + 1
                                    curr_num = num_counters[lvl]
                                    clean_p = re.sub(r"^\d+\.\s+", "", p_text)
                                    para_lines.append(f"{indent}{curr_num}. {clean_p}")
                                else:
                                    num_counters[lvl] = 0
                                    clean_p = re.sub(r"^•\s*", "", p_text)
                                    para_lines.append(f"{indent}- {clean_p}")

                        if para_lines:
                            slide_blocks.append("\n".join(para_lines))
                        continue

                    # Fallback for unknown shape types with text/content
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_blocks.append(shape.text.strip())

                # 7. Slide Notes Extraction (preserving bullets & numbering)
                if slide.has_notes_slide:
                    try:
                        notes_tf = slide.notes_slide.notes_text_frame
                        if notes_tf and notes_tf.text and notes_tf.text.strip():
                            notes_lines = []
                            num_counters_notes = {}
                            for p in notes_tf.paragraphs:
                                p_str = p.text.strip()
                                if not p_str:
                                    continue
                                lvl = getattr(p, "level", 0) or 0
                                indent = "  " * lvl
                                is_num = self._is_numbered_paragraph(p) or bool(re.match(r"^\d+\.\s+", p_str))
                                if is_num:
                                    num_counters_notes[lvl] = num_counters_notes.get(lvl, 0) + 1
                                    clean_p = re.sub(r"^\d+\.\s+", "", p_str)
                                    notes_lines.append(f"{indent}{num_counters_notes[lvl]}. {clean_p}")
                                else:
                                    num_counters_notes[lvl] = 0
                                    clean_p = re.sub(r"^•\s*", "", p_str)
                                    if clean_p.startswith("- "):
                                        notes_lines.append(f"{indent}{clean_p}")
                                    else:
                                        notes_lines.append(f"{indent}- {clean_p}")

                            if notes_lines:
                                slide_blocks.append("### Notes:\n" + "\n".join(notes_lines))
                    except Exception as ex:
                        print(f"[DEBUG] Slide notes extraction error: {ex}")

                slide_content = "\n\n".join(slide_blocks).strip()
                if not slide_content:
                    slide_content = f"*(Empty Slide {slide_idx})*"
                slide_markdowns.append(slide_content)

            # Join slides strictly with '\n\n---\n\n'
            return "\n\n---\n\n".join(slide_markdowns)

        except Exception as e:
            return f"⚠️ PPTX Conversion Error: {str(e)}"

    def save_from_markdown(self, markdown_content: str, out_path: str) -> str:
        """
        Converts Markdown text into a PowerPoint presentation (.pptx).
        Parses headings, bullet levels, bold/italic styles, tables, charts, images, and slide notes.
        Splits slides strictly by '\n\n---\n\n' markers, or automatically segments by major section headings.
        Note: When a long bullet list is chunked across slides, sub-bullets may start at the top of a new slide without their parent header bullet (known aesthetic trade-off for auto-paginated slides).
        """
        try:
            from src.core.converters import prepare_markdown_for_export, parse_table_rows
            markdown_content = prepare_markdown_for_export(markdown_content)

            import pptx
            from pptx.util import Inches, Pt
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
            try:
                from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
            except ImportError:
                MSO_ANCHOR = None
                MSO_AUTO_SIZE = None
                PP_ALIGN = None
            from src.services.media_asset_manager import MediaAssetManager
            from src.ui_flet.helpers.image_token_helper import find_all_image_tokens

            asset_mgr = MediaAssetManager()
            prs = pptx.Presentation()
            # Set default 16:9 Widescreen aspect ratio (13.333 in x 7.5 in)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]  # Blank slide layout

            slide_blocks = self._split_markdown_into_slide_blocks(markdown_content)

            if not slide_blocks:
                slide_blocks = [markdown_content.strip()]

            for block in slide_blocks:
                # 1. Extract Slide Notes if present
                notes_text = None
                notes_match = re.search(r"(?:^|\n)### Notes:\s*\n(.*)$|(?:^|\n)Notes:\s*\n(.*)$", block, re.DOTALL)
                if notes_match:
                    notes_text = (notes_match.group(1) if notes_match.group(1) is not None else notes_match.group(2) or "").strip()
                    block = block[:notes_match.start()].strip()

                lines = block.split("\n")

                title_text = None
                content_lines = []
                in_block_code = False

                for idx, l in enumerate(lines):
                    stripped = l.strip()
                    if stripped.startswith("```"):
                        in_block_code = not in_block_code
                        content_lines.append(l)
                        continue
                    if not title_text and not in_block_code and idx < 3 and (stripped.startswith("# ") or stripped.startswith("## ") or stripped.startswith("### ")):
                        title_text = re.sub(r"^#+\s*", "", stripped)
                    elif re.match(r"^\*\s*\(?Empty Slide.*\)?\s*\*$", stripped, re.IGNORECASE):
                        continue
                    else:
                        content_lines.append(l)

                if title_text:
                    title_layout = prs.slide_layouts[5]  # Title Only layout
                    slide = prs.slides.add_slide(title_layout)
                    initial_offset, actual_title_h, title_font_sz = self._calc_title_geometry(title_text, Inches)
                    if hasattr(slide.shapes, "title") and slide.shapes.title:
                        title_shape = slide.shapes.title
                        title_shape.text = title_text
                        try:
                            title_shape.left = Inches(0.8)
                            title_shape.top = Inches(0.5)
                            title_shape.width = Inches(11.7)
                            title_shape.height = actual_title_h
                            tf_t = title_shape.text_frame
                            tf_t.word_wrap = True
                            tf_t.margin_top = Inches(0.0)
                            tf_t.margin_bottom = Inches(0.0)
                            tf_t.margin_left = Inches(0.0)
                            tf_t.margin_right = Inches(0.0)
                            for p in tf_t.paragraphs:
                                if PP_ALIGN is not None:
                                    try:
                                        p.alignment = PP_ALIGN.LEFT
                                    except Exception:
                                        pass
                                for r in p.runs:
                                    r.font.size = Pt(title_font_sz)
                        except Exception:
                            pass
                    top_offset = initial_offset
                    start_top_offset = initial_offset
                else:
                    slide = prs.slides.add_slide(blank_layout)
                    top_offset = Inches(0.8)
                    start_top_offset = Inches(0.8)

                # Attach Slide Notes if present
                if notes_text:
                    try:
                        slide.notes_slide.notes_text_frame.text = notes_text
                    except Exception as ex:
                        print(f"[DEBUG] Failed to set slide notes: {ex}")

                # Overflow protection threshold
                max_slide_top = prs.slide_height - Inches(0.8)

                # Process content lines (detecting tables, charts, images, bullet lists, paragraphs)
                i = 0
                pending_chart_type = None
                pending_chart_title = None
                pending_chart_legend = None

                while i < len(content_lines):
                    line = content_lines[i]
                    stripped = line.strip()

                    if not stripped:
                        i += 1
                        continue

                    # 1.5 Chart footprint comment detection <!-- chart: type title="..." legend="..." -->
                    chart_comment_match = re.match(
                        r"^<!--\s*chart:\s*(\w+)(?:.*?title=[\"'](.*?)[\"'])?(?:.*?legend=[\"'](.*?)[\"'])?\s*-->$",
                        stripped,
                        re.IGNORECASE
                    )
                    if chart_comment_match:
                        pending_chart_type = chart_comment_match.group(1).lower()
                        pending_chart_title = chart_comment_match.group(2) or ""
                        pending_chart_legend = chart_comment_match.group(3) or ""
                        i += 1
                        continue

                    # 2. Markdown Table or Chart Detection (| col1 | col2 |)
                    if stripped.startswith("|") and stripped.endswith("|"):
                        table_block = []
                        while i < len(content_lines) and content_lines[i].strip().startswith("|") and content_lines[i].strip().endswith("|"):
                            table_block.append(content_lines[i].strip())
                            i += 1

                        # Use parse_table_rows helper from converters.py
                        rows_data = parse_table_rows(table_block)
                        if rows_data:
                            # If chart footprint comment was present, construct native PowerPoint Chart!
                            if pending_chart_type and len(rows_data) > 1:
                                try:
                                    slide, top_offset, start_top_offset = self._ensure_space(
                                        slide, top_offset, start_top_offset, max_slide_top, prs,
                                        title_text, notes_text, blank_layout, Inches, Pt, PP_ALIGN, Inches(4.8)
                                    )
                                    chart_data = CategoryChartData()
                                    headers = rows_data[0]
                                    categories = [r[0] for r in rows_data[1:]]
                                    chart_data.categories = categories

                                    for col_idx in range(1, len(headers)):
                                        series_name = headers[col_idx]
                                        series_vals = []
                                        for r in rows_data[1:]:
                                            try:
                                                v = float(r[col_idx]) if col_idx < len(r) else 0.0
                                            except (ValueError, TypeError):
                                                v = 0.0
                                            series_vals.append(v)
                                        chart_data.add_series(series_name, series_vals)

                                    chart_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
                                    if pending_chart_type == "bar":
                                        chart_enum = XL_CHART_TYPE.BAR_CLUSTERED
                                    elif pending_chart_type == "pie":
                                        chart_enum = XL_CHART_TYPE.PIE
                                    elif pending_chart_type == "line":
                                        chart_enum = XL_CHART_TYPE.LINE

                                    chart_shape = slide.shapes.add_chart(chart_enum, Inches(0.8), top_offset, Inches(11.7), Inches(4.5), chart_data)
                                    chart_obj = chart_shape.chart

                                    if pending_chart_title:
                                        chart_obj.has_title = True
                                        if hasattr(chart_obj.chart_title, "text_frame") and chart_obj.chart_title.text_frame.paragraphs:
                                            chart_obj.chart_title.text_frame.paragraphs[0].text = pending_chart_title
                                    else:
                                        chart_obj.has_title = False

                                    chart_obj.has_legend = True
                                    try:
                                        chart_obj.legend.include_in_layout = False
                                    except Exception:
                                        pass
                                    if pending_chart_legend:
                                        l_str = pending_chart_legend.lower()
                                        if "right" in l_str:
                                            chart_obj.legend.position = XL_LEGEND_POSITION.RIGHT
                                        elif "top" in l_str:
                                            chart_obj.legend.position = XL_LEGEND_POSITION.TOP
                                        elif "left" in l_str:
                                            chart_obj.legend.position = XL_LEGEND_POSITION.LEFT
                                        elif "bottom" in l_str:
                                            chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM

                                    top_offset += Inches(4.8)
                                    pending_chart_type = None
                                    pending_chart_title = None
                                    pending_chart_legend = None
                                    continue
                                except Exception as chart_ex:
                                    print(f"[DEBUG] Failed to create chart from footprint: {chart_ex}")

                            # Native PowerPoint Table with Continuation & Cumulative Height Slicing
                            header_row = rows_data[0]
                            data_rows = rows_data[1:]
                            num_cols = max(len(r) for r in rows_data) if rows_data else 1
                            col_char_capacities = [max(20, int(75 / num_cols))] * num_cols

                            rendered_once = False
                            while data_rows or not rendered_once:
                                avail = max_slide_top - top_offset - Inches(0.4)
                                header_h = self._calc_row_h(header_row, col_char_capacities, Inches)

                                slice_data_rows = []
                                used_h = header_h

                                while data_rows:
                                    next_r_h = self._calc_row_h(data_rows[0], col_char_capacities, Inches)
                                    if used_h + next_r_h > avail and slice_data_rows:
                                        break
                                    slice_data_rows.append(data_rows.pop(0))
                                    used_h += next_r_h

                                table_slice = [header_row] + slice_data_rows
                                row_count = len(table_slice)
                                col_count = num_cols

                                slide, top_offset, start_top_offset = self._ensure_space(
                                    slide, top_offset, start_top_offset, max_slide_top, prs,
                                    title_text, notes_text, blank_layout, Inches, Pt, PP_ALIGN, used_h + Inches(0.2)
                                )

                                table_shape = slide.shapes.add_table(row_count, col_count, Inches(0.8), top_offset, Inches(11.7), used_h)
                                table = table_shape.table
                                for r_idx, row in enumerate(table_slice):
                                    for c_idx, val in enumerate(row):
                                        if c_idx < col_count:
                                            cell = table.cell(r_idx, c_idx)
                                            if MSO_ANCHOR is not None:
                                                try:
                                                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                                                except Exception:
                                                    pass
                                            cell.text = ""  # Clear default text
                                            tf = cell.text_frame
                                            tf.word_wrap = True
                                            cell_text = val.replace("<br>", "\n")
                                            lines_in_cell = cell_text.split("\n")
                                            for l_idx, c_line in enumerate(lines_in_cell):
                                                p = tf.paragraphs[0] if l_idx == 0 else tf.add_paragraph()
                                                if r_idx == 0:
                                                    self._parse_formatted_runs(p, c_line, Pt, default_size=15, force_bold=True)
                                                else:
                                                    self._parse_formatted_runs(p, c_line, Pt, default_size=15)

                                top_offset += used_h + Inches(0.5)
                                rendered_once = True
                                if data_rows:
                                    slide, top_offset, start_top_offset = self._ensure_space(
                                        slide, top_offset, start_top_offset, max_slide_top, prs,
                                        title_text, notes_text, blank_layout, Inches, Pt, PP_ALIGN, Inches(2.0)
                                    )
                            pending_chart_type = None
                        continue

                    # 3. Image Detection (Markdown ![alt](path) and HTML <img...> / <p align="...">)
                    img_tokens = find_all_image_tokens(stripped) if ("!" in stripped or "<img" in stripped or "<p" in stripped or "<div" in stripped or "<center" in stripped) else []
                    if img_tokens:
                        tok = img_tokens[0]
                        img_uri = tok.src
                        out_dir = os.path.dirname(out_path) if out_path else None
                        resolved = asset_mgr.resolve_uri(img_uri, base_dir=out_dir)
                        if not os.path.exists(resolved) and img_uri and os.path.exists(img_uri):
                            resolved = img_uri

                        if os.path.exists(resolved):
                            try:
                                slide, top_offset, start_top_offset = self._ensure_space(
                                    slide, top_offset, start_top_offset, max_slide_top, prs,
                                    title_text, notes_text, blank_layout, Inches, Pt, PP_ALIGN, Inches(3.8)
                                )
                                # Calculate picture width based on token width if specified
                                pic_w = Inches(7.0)
                                if tok.width:
                                    if tok.width.endswith("%"):
                                        try:
                                            pct = float(tok.width.rstrip("%"))
                                            pic_w = Inches(min(8.0, max(1.0, 8.0 * (pct / 100.0))))
                                        except Exception:
                                            pic_w = Inches(7.0)
                                    else:
                                        try:
                                            px_val = float(tok.width.rstrip("px"))
                                            pic_w = Inches(min(8.0, max(1.0, px_val / 96.0)))
                                        except Exception:
                                            pic_w = Inches(7.0)

                                # Calculate horizontal alignment
                                slide_w = getattr(prs, "slide_width", Inches(13.333))
                                if tok.align == "center":
                                    left_pos = max(Inches(0.5), (slide_w - pic_w) / 2)
                                elif tok.align == "right":
                                    left_pos = max(Inches(0.5), slide_w - pic_w - Inches(0.8))
                                else:
                                    left_pos = Inches(0.8)

                                slide.shapes.add_picture(resolved, left_pos, top_offset, width=pic_w)
                                top_offset += Inches(3.8)
                            except Exception as ex:
                                print(f"[DEBUG] Failed to insert picture in PPTX export: {ex}")
                        i += 1
                        continue

                    # 4. Text Frame (Bullet lists & formatted paragraphs with pagination)
                    # Sub-heading + Table cohesion check
                    is_sub_heading = bool(re.match(r"^#{1,6}\s+", stripped)) or (stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4) or bool(re.match(r"^\*\*Mục\s+\d+", stripped, re.IGNORECASE))
                    if is_sub_heading and i + 1 < len(content_lines) and content_lines[i + 1].strip().startswith("|"):
                        slide, top_offset, start_top_offset = self._ensure_space(
                            slide, top_offset, start_top_offset, max_slide_top, prs,
                            title_text, notes_text, blank_layout, Inches, Pt, PP_ALIGN, Inches(2.8)
                        )

                    tx_lines = []
                    while i < len(content_lines):
                        l_str = content_lines[i].strip()
                        if l_str.startswith("```"):
                            tx_lines.append(content_lines[i])
                            i += 1
                            continue
                        if not l_str or (l_str.startswith("|") and l_str.endswith("|")) or re.match(r"^!\[([^\]]*)\]\((.+?\.(?:png|jpg|jpeg|gif|svg|webp|bmp|ico)|https?://\S+|@media/\S+?|[^\n)]+)\)$", l_str, re.IGNORECASE):
                            break
                        tx_lines.append(content_lines[i])
                        i += 1

                    while tx_lines:
                        first_str = tx_lines[0].strip() if tx_lines else ""
                        is_heading_start = bool(re.match(r"^#{1,6}\s+", first_str)) or (first_str.startswith("**") and first_str.endswith("**") and len(first_str) > 4) or bool(re.match(r"^\*\*Mục\s+\d+", first_str, re.IGNORECASE))

                        if (top_offset > max_slide_top - Inches(1.2) or (is_heading_start and max_slide_top - top_offset < Inches(2.2))) and top_offset > start_top_offset:
                            slide, top_offset, start_top_offset = self._create_continuation_slide(
                                prs, title_text, notes_text, blank_layout, Inches, Pt, PP_ALIGN
                            )

                        avail_height = max_slide_top - top_offset
                        line_height = Inches(0.24)
                        max_lines_fit = max(1, int((avail_height - Inches(0.1)) / line_height))

                        chunk = []
                        accumulated_visual_lines = 0
                        for bl in tx_lines:
                            leading_spaces = len(bl) - len(bl.lstrip(" "))
                            line_lvl = min(8, leading_spaces // 2)
                            v_lines = self._get_visual_line_count(bl, line_lvl)

                            if accumulated_visual_lines + v_lines > max_lines_fit and chunk:
                                break
                            chunk.append(bl)
                            accumulated_visual_lines += v_lines

                        # Atomic Unit Grouping
                        if len(chunk) > 1:
                            bold_indices = []
                            c_in_code = False
                            for idx, item in enumerate(chunk):
                                s_item = item.strip()
                                if s_item.startswith("```"):
                                    c_in_code = not c_in_code
                                    continue
                                is_list_item = bool(re.match(r"^(?:[*\-+]|\d+\.)\s+", s_item))
                                is_standalone_subheader = not is_list_item and (re.match(r"^#{1,6}\s+", s_item) or (s_item.startswith("**") and s_item.endswith("**") and len(s_item) > 4))
                                if not c_in_code and is_standalone_subheader:
                                    bold_indices.append(idx)

                            if bold_indices and bold_indices[-1] > 0:
                                last_bold_idx = bold_indices[-1]
                                chunk = chunk[:last_bold_idx]
                                accumulated_visual_lines = sum(self._get_visual_line_count(item, min(8, (len(item) - len(item.lstrip(" "))) // 2)) for item in chunk)

                        # Intro Line + Code Block Cohesion Rewind
                        if len(chunk) > 0 and len(chunk) < len(tx_lines):
                            last_item = chunk[-1].strip()
                            next_item = tx_lines[len(chunk)].strip()
                            if (last_item.endswith(":") or last_item.endswith(":**")) and next_item.startswith("```"):
                                chunk.pop()
                                accumulated_visual_lines = sum(self._get_visual_line_count(item, min(8, (len(item) - len(item.lstrip(" "))) // 2)) for item in chunk)

                        if not chunk and tx_lines:
                            first_bl = tx_lines[0]
                            chunk = [first_bl]
                            leading_spaces = len(first_bl) - len(first_bl.lstrip(" "))
                            accumulated_visual_lines = self._get_visual_line_count(first_bl, min(8, leading_spaces // 2))

                        tx_lines = tx_lines[len(chunk):]

                        if accumulated_visual_lines <= 8:
                            base_font_sz = 18
                        else:
                            base_font_sz = 16

                        box_height = Inches(0.28 * accumulated_visual_lines + 0.15)
                        txBox = slide.shapes.add_textbox(Inches(0.8), top_offset, Inches(11.7), box_height)
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        try:
                            tf.margin_top = Inches(0.0)
                            tf.margin_bottom = Inches(0.0)
                            tf.margin_left = Inches(0.0)
                            tf.margin_right = Inches(0.0)
                        except Exception:
                            pass

                        if MSO_AUTO_SIZE is not None:
                            try:
                                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            except Exception:
                                pass

                        first_para = True
                        in_code_block = False
                        for bl in chunk:
                            bl_str = bl.strip()
                            if not bl_str:
                                continue

                            if bl_str.startswith("```"):
                                in_code_block = not in_code_block
                                continue

                            if first_para:
                                p = tf.paragraphs[0]
                                first_para = False
                            else:
                                p = tf.add_paragraph()

                            leading_spaces = len(bl) - len(bl.lstrip(" "))
                            level = min(8, leading_spaces // 2)
                            p.level = level

                            if in_code_block:
                                p.level = 0
                                clean_code_line = re.sub(r"^[*\-+]\s+", "", bl_str)
                                try:
                                    p.space_before = Pt(2)
                                    p.space_after = Pt(2)
                                except Exception:
                                    pass
                                self._parse_formatted_runs(p, f"`{clean_code_line}`", Pt, default_size=16)
                                continue

                            h_match = re.match(r"^(#{1,6})\s+(.*)$", bl_str) if (leading_spaces <= 2 and not in_code_block) else None
                            if h_match:
                                h_lvl = len(h_match.group(1))
                                clean_text = h_match.group(2)
                                if h_lvl == 1:
                                    h_size = 28
                                elif h_lvl == 2:
                                    h_size = 25
                                elif h_lvl == 3:
                                    h_size = 22
                                else:
                                    h_size = 20

                                p.level = 0
                                try:
                                    p.space_before = Pt(14)
                                    p.space_after = Pt(6)
                                except Exception:
                                    pass
                                self._parse_formatted_runs(p, clean_text, Pt, default_size=h_size, force_bold=True)
                                continue

                            is_bold_subheader = (bl_str.startswith("**") and bl_str.endswith("**") and len(bl_str) > 4) or bool(re.match(r"^\*\*Mục\s+\d+", bl_str, re.IGNORECASE))
                            if is_bold_subheader or (bl_str.startswith("**") and "**" in bl_str[2:] and not bl_str.startswith("- ") and not bl_str.startswith("* ") and not bl_str.startswith("• ")):
                                p.level = 0
                                try:
                                    p.space_before = Pt(12)
                                    p.space_after = Pt(4)
                                except Exception:
                                    pass
                                self._parse_formatted_runs(p, bl_str, Pt, default_size=max(20, base_font_sz + 3), force_bold=True)
                                continue

                            is_numbered = bool(re.match(r"^\d+\.\s+", bl_str))
                            is_bullet = bool(re.match(r"^[*\-+]\s+", bl_str))

                            if is_numbered:
                                clean_text = bl_str
                            elif is_bullet:
                                stripped_text = re.sub(r"^[*\-+]\s+", "", bl_str)
                                clean_text = f"•  {stripped_text}"
                            elif level > 0:
                                clean_text = f"•  {bl_str}"
                            else:
                                clean_text = bl_str

                            font_sz = max(16, base_font_sz - (level if base_font_sz > 16 else 0))
                            try:
                                p.space_before = Pt(2)
                                p.space_after = Pt(8)
                            except Exception:
                                pass
                            self._parse_formatted_runs(p, clean_text, Pt, default_size=font_sz)

                        top_offset += box_height

            out_dir = os.path.dirname(out_path)
            if out_dir:
                os.makedirs(out_dir, exist_ok=True)

            prs.save(out_path)
            return f"Saved to PowerPoint -> {os.path.basename(out_path)}"

        except Exception as e:
            return f"⚠️ Export Error: Failed to save PPTX document — {str(e)}"


# Automatically register module when imported
ModuleRegistry.register(PPTXModule())
