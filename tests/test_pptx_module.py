"""
Unit & Integration Tests for PowerPoint (PPTX) Document Module.
"""
import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.abspath("."))


class TestPPTXModule(unittest.TestCase):
    def test_pptx_module_registration(self):
        from src.core.registry import ModuleRegistry
        from src.modules.pptx_module import PPTXModule

        module = ModuleRegistry.get_module_by_extension(".pptx")
        self.assertIsNotNone(module)
        self.assertIsInstance(module, PPTXModule)
        self.assertEqual(module.name, "PowerPoint")
        self.assertIn(".pptx", module.file_extensions)

    def test_check_dependencies_mapping(self):
        from src.modules.pptx_module import PPTXModule

        module = PPTXModule()
        self.assertIn("python-pptx", module.required_dependencies)
        missing = module.check_dependencies()
        # python-pptx should be mapped to import pptx
        # If python-pptx is installed, missing should not include python-pptx
        try:
            import pptx
            self.assertNotIn("python-pptx", missing)
        except ImportError:
            self.assertIn("python-pptx", missing)

    def test_constants_registration(self):
        from src.ui_flet.constants import MODES, MODE_DISPLAY_KEYS, IN_FILETYPES, OUT_FILETYPES

        self.assertIn("MD -> PowerPoint", MODES)
        self.assertIn("PowerPoint -> MD", MODES)
        self.assertIn("MD -> PowerPoint", MODE_DISPLAY_KEYS)
        self.assertIn("PowerPoint -> MD", MODE_DISPLAY_KEYS)
        self.assertIn(".pptx", IN_FILETYPES)
        self.assertIn(".pptx", OUT_FILETYPES)

    def test_save_and_load_roundtrip(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        md_input = "## Slide 1 Title\n- Bullet item 1\n  - Sub bullet 2\n\n---\n\n## Slide 2 Title\n- Item A\n- Item B"

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_output.pptx")

            # Test save_from_markdown
            res = module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))
            self.assertIn("Saved to PowerPoint", res)

            # Test load_to_markdown
            extracted_md = module.load_to_markdown(out_pptx)
            self.assertIn("## Slide 1 Title", extracted_md)
            self.assertIn("## Slide 2 Title", extracted_md)
            self.assertIn("---", extracted_md)
            self.assertIn("Bullet item 1", extracted_md)

    def test_table_and_notes_export(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        md_input = (
            "## Table Slide\n"
            "| Header 1 | Header 2 |\n"
            "| --- | --- |\n"
            "| Val 1 | Val 2 |\n\n"
            "### Notes:\n"
            "Test note content"
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_table_notes.pptx")
            module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))

            # Inspect created PPTX presentation
            prs = pptx.Presentation(out_pptx)
            slide = prs.slides[0]
            
            # Check table existence
            has_table = any(s.has_table for s in slide.shapes)
            self.assertTrue(has_table)

            # Check slide notes
            self.assertIn("Test note content", slide.notes_slide.notes_text_frame.text)

    def test_chart_footprint_roundtrip(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        md_input = (
            "## Chart Slide\n"
            "<!-- chart: column -->\n"
            "| Category | Series 1 | Series 2 |\n"
            "| --- | --- | --- |\n"
            "| Cat A | 10.5 | 20.0 |\n"
            "| Cat B | 15.0 | 25.5 |"
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_chart.pptx")
            module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))

            prs = pptx.Presentation(out_pptx)
            slide = prs.slides[0]
            has_chart = any(s.has_chart for s in slide.shapes)
            self.assertTrue(has_chart)

    def test_numbered_list_roundtrip(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        md_input = (
            "## Numbering Slide\n"
            "1. First item\n"
            "2. Second item"
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_num.pptx")
            module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))

            # Extract back to Markdown
            extracted = module.load_to_markdown(out_pptx)
            self.assertIn("1. First item", extracted)
            self.assertIn("2. Second item", extracted)

    def test_table_roundtrip(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        md_input = (
            "## Table Slide\n"
            "| Left | Center | Right |\n"
            "| --- | --- | --- |\n"
            "| L1 | C1 | R1 |"
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_table.pptx")
            module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))

            # Extract back to Markdown
            extracted = module.load_to_markdown(out_pptx)
            self.assertIn("| --- | --- | --- |", extracted)

    def test_long_text_overflow_pagination(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        # 30 lines of long text in a single block without title
        long_lines = [f"- Paragraph item {i}: " + ("Sample long text " * 5) for i in range(1, 31)]
        md_input = "\n".join(long_lines)

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_overflow.pptx")
            module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))

            prs = pptx.Presentation(out_pptx)
            # 30 lines should be split into multiple slides (>1 slide)
            self.assertGreater(len(prs.slides), 1)
            # Ensure no slide is empty
            for s_idx, slide in enumerate(prs.slides):
                self.assertGreater(len(slide.shapes), 0, f"Slide {s_idx+1} is empty!")

    def test_long_paragraph_chunking(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        # Markdown with very long paragraphs (each >300 chars) that wrap visually into multiple lines
        p1 = "Paragraph 1: " + ("This is a long sentence with detailed information about topic A. " * 8)
        p2 = "Paragraph 2: " + ("This is another long sentence detailing complex systems and workflows. " * 8)
        p3 = "Paragraph 3: " + ("Final paragraph discussing architectural trade-offs and future developments. " * 8)
        md_input = f"## Big Content Slide\n\n{p1}\n\n{p2}\n\n{p3}"

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_long_paragraphs.pptx")
            module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))

            prs = pptx.Presentation(out_pptx)
            # Long paragraphs should wrap visually and split into multiple slides when exceeding height
            self.assertGreaterEqual(len(prs.slides), 1)
            extracted = module.load_to_markdown(out_pptx)
            self.assertIn("Paragraph 1", extracted)
            self.assertIn("Paragraph 2", extracted)
            self.assertIn("Paragraph 3", extracted)

    def test_header_only_table_infinite_loop_protection(self):
        try:
            import pptx
        except ImportError:
            self.skipTest("python-pptx not installed")

        from src.modules.pptx_module import PPTXModule
        module = PPTXModule()

        # Table with header row and separator row, but 0 data rows
        md_input = "## Header Only Table\n| Col1 | Col2 |\n| --- | --- |\n"

        with tempfile.TemporaryDirectory() as tmp_dir:
            out_pptx = os.path.join(tmp_dir, "test_header_only_table.pptx")
            # Should complete without infinite loop hanging
            res = module.save_from_markdown(md_input, out_pptx)
            self.assertTrue(os.path.exists(out_pptx))
            self.assertIn("Saved to PowerPoint", res)

            prs = pptx.Presentation(out_pptx)
            self.assertEqual(len(prs.slides), 1)


if __name__ == "__main__":
    unittest.main()

